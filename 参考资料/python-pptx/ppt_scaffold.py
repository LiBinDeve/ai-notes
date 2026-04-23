from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# =========================
# 第三层：样式配置
# =========================


@dataclass
class Theme:
    # 画布尺寸（16:9）
    slide_width: float = 13.333
    slide_height: float = 7.5

    # 页面边距
    margin_left: float = 0.6
    margin_top: float = 0.4
    content_width: float = 12.1

    # 字体
    font_name: str = "Microsoft YaHei"
    title_size: int = 24
    subtitle_size: int = 14
    question_size: int = 28
    body_size: int = 20
    small_size: int = 12

    # 颜色
    color_text: tuple[int, int, int] = (34, 34, 34)
    color_muted: tuple[int, int, int] = (110, 110, 110)
    color_primary: tuple[int, int, int] = (44, 92, 197)
    color_bg: tuple[int, int, int] = (255, 255, 255)
    color_box: tuple[int, int, int] = (245, 247, 250)


# =========================
# 基础工具
# =========================


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def new_presentation(theme: Theme | None = None) -> Presentation:
    theme = theme or Theme()
    prs = Presentation()
    prs.slide_width = Inches(theme.slide_width)
    prs.slide_height = Inches(theme.slide_height)
    return prs


# =========================
# 第二层：组件函数
# =========================


def add_blank_slide(prs: Presentation):
    """添加空白页。"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_name: str = "Microsoft YaHei",
    font_size: int = 20,
    bold: bool = False,
    color: tuple[int, int, int] = (34, 34, 34),
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
    margin_left: float = 0.05,
    margin_right: float = 0.05,
    margin_top: float = 0.03,
    margin_bottom: float = 0.03,
    fill_color: tuple[int, int, int] | None = None,
    line_color: tuple[int, int, int] | None = None,
    line_width_pt: float = 1.0,
):
    """添加文本框。坐标单位统一用英寸，后面改布局最方便。"""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))

    # 文本框填充
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)

    # 文本框边框
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width_pt)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)

    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill_color: tuple[int, int, int] = (245, 247, 250),
    line_color: tuple[int, int, int] | None = None,
    line_width_pt: float = 1.0,
):
    """添加矩形，可作为底板、标签背景、分区卡片。"""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)

    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width_pt)

    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: tuple[int, int, int] = (44, 92, 197),
    width_pt: float = 1.5,
):
    """添加直线。"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width_pt)
    return line


def add_image(slide, image_path: str | Path, x: float, y: float, w: float | None = None, h: float | None = None):
    """添加图片。"""
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), **kwargs)


# =========================
# 第一层：页面函数
# =========================


def add_cover_slide(prs: Presentation, title: str, subtitle: str = "", theme: Theme | None = None):
    """封面页。"""
    theme = theme or Theme()
    slide = add_blank_slide(prs)

    # 顶部强调线
    add_line(slide, 0.7, 0.8, 4.8, 0.8, color=theme.color_primary, width_pt=3)

    add_textbox(
        slide, title,
        0.8, 1.1, 11.2, 1.0,
        font_name=theme.font_name,
        font_size=34,
        bold=True,
        color=theme.color_text,
    )

    if subtitle:
        add_textbox(
            slide, subtitle,
            0.8, 2.1, 10.5, 0.6,
            font_name=theme.font_name,
            font_size=theme.subtitle_size,
            color=theme.color_muted,
        )

    return slide


def add_question_slide(
    prs: Presentation,
    page_title: str,
    question: str,
    options: Sequence[str],
    theme: Theme | None = None,
):
    """题目页。"""
    theme = theme or Theme()
    slide = add_blank_slide(prs)

    # 页标题
    add_textbox(
        slide, page_title,
        theme.margin_left, 0.35, 8.5, 0.45,
        font_name=theme.font_name,
        font_size=theme.title_size,
        bold=True,
        color=theme.color_primary,
    )

    # 题干底板
    add_rect(
        slide,
        0.65, 0.95, 12.0, 1.15,
        fill_color=theme.color_box,
    )
    add_textbox(
        slide, question,
        0.85, 1.15, 11.5, 0.8,
        font_name=theme.font_name,
        font_size=theme.question_size,
        bold=True,
        color=theme.color_text,
    )

    # 选项区
    y = 2.45
    for opt in options:
        add_rect(slide, 0.85, y - 0.05, 11.2, 0.58, fill_color=(255, 255, 255), line_color=(220, 224, 230))
        add_textbox(
            slide, opt,
            1.05, y, 10.8, 0.4,
            font_name=theme.font_name,
            font_size=theme.body_size,
            color=theme.color_text,
        )
        y += 0.75

    return slide


def add_summary_slide(
    prs: Presentation,
    title: str,
    bullets: Iterable[str],
    theme: Theme | None = None,
):
    """总结页。"""
    theme = theme or Theme()
    slide = add_blank_slide(prs)

    add_textbox(
        slide, title,
        0.8, 0.45, 10.0, 0.55,
        font_name=theme.font_name,
        font_size=theme.title_size + 4,
        bold=True,
        color=theme.color_text,
    )
    add_line(slide, 0.8, 1.0, 3.2, 1.0, color=theme.color_primary, width_pt=2.5)

    y = 1.45
    for item in bullets:
        add_textbox(
            slide, f"• {item}",
            1.0, y, 11.2, 0.55,
            font_name=theme.font_name,
            font_size=theme.body_size,
            color=theme.color_text,
        )
        y += 0.65

    return slide


# =========================
# 示例：如何调用
# =========================


def build_demo_ppt(output_path: str = "demo_ppt_template.pptx") -> None:
    theme = Theme(
        font_name="Microsoft YaHei",
        title_size=24,
        question_size=28,
        body_size=20,
    )

    prs = new_presentation(theme)

    add_cover_slide(
        prs,
        title="四中全会题库",
        subtitle="自动生成示例｜你以后只需要改数据和页面函数",
        theme=theme,
    )

    add_question_slide(
        prs,
        page_title="（2026年中央经济工作会议）",
        question="1. 下列哪项表述正确？",
        options=[
            "A. 选项一",
            "B. 选项二",
            "C. 选项三",
            "D. 选项四",
        ],
        theme=theme,
    )

    add_summary_slide(
        prs,
        title="使用说明",
        bullets=[
            "改样式：只改 Theme 里的字体、字号、颜色、边距。",
            "改组件：只改 add_textbox / add_rect / add_image / add_line。",
            "改页面：新增 add_xxx_slide 页面函数即可。",
            "接 YAML / Excel / JSON 数据时，只需要在外面读取后传进来。",
        ],
        theme=theme,
    )

    prs.save(output_path)
    print(f"已生成：{output_path}")


if __name__ == "__main__":
    build_demo_ppt()
